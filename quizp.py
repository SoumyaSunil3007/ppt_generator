from pptx import Presentation
from pptx.util import Inches
import random

# Sample questions and answers
questions_and_answers = [
    ("What is the capital of France?", "Paris", "paris.jpg"),
    ("Who wrote 'To Kill a Mockingbird'?", "Harper Lee", "harper_lee.jpg"),
    ("What is the chemical symbol for water?", "H2O", "water.jpg")
]

# Function to create a slide with a question, answer, and image
def add_slide(prs, question, answer, image_path):
    slide_layout = prs.slide_layouts[5]  # Assuming we are using a title and content layout
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Set the title to the question
    title_shape = shapes.title
    title_shape.text = "Question: " + question

    # Add the image
    left_inch = Inches(5)
    top_inch = Inches(2)
    width_inch = Inches(3)
    height_inch = Inches(3)
    pic = slide.shapes.add_picture(image_path, left_inch, top_inch, width=width_inch, height=height_inch)

    # Add the answer as a text box
    left_inch = Inches(1)
    top_inch = Inches(2)
    width_inch = Inches(8)
    height_inch = Inches(1.5)
    txBox = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Answer: " + answer

# Create a presentation object
prs = Presentation("template.pptx")  # Load the template

# Identify the layout index for the "Lilac" slide layout
# You need to check the index of "Lilac" layout in your template
lilac_layout_idx = 0

# Add the "Lilac" slide as the first slide in the presentation
lilac_slide_layout = prs.slide_layouts[lilac_layout_idx]
lilac_slide = prs.slides.add_slide(lilac_slide_layout)
title_shape = lilac_slide.shapes.title
title_shape.text = "Lilac"  # Title of the slide

# Shuffle the questions and answers
random.shuffle(questions_and_answers)

# Add slides for each question-answer pair
for question, answer, image_path in questions_and_answers:
    add_slide(prs, question, answer, image_path)

# Save the presentation
prs.save("random_quiz_with_lilac.pptx")
